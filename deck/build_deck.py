"""Generate the Build 2026 Recap slide deck (PowerPoint .pptx) with python-pptx.

    python deck/build_deck.py

Produces: deck/Build2026-Recap-Agents.pptx  (16:9, ~20 slides, with speaker notes).

The content mirrors the repo and the DEMO_SCRIPT runbook, so the deck and the code
stay in sync. Re-run any time you tweak the story.
"""

from __future__ import annotations

import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

# ── palette ───────────────────────────────────────────────────────────────────
AZURE = RGBColor(0x00, 0x78, 0xD4)
DARK = RGBColor(0x20, 0x20, 0x20)
GREY = RGBColor(0x60, 0x60, 0x60)
LIGHT = RGBColor(0xF3, 0xF6, 0xFB)
CODEBG = RGBColor(0x1E, 0x1E, 0x1E)
CODEFG = RGBColor(0xE6, 0xE6, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 16:9
W = Emu(12192000)
H = Emu(6858000)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def _tb(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _band(slide, color=AZURE, height=Emu(120000), top=Emu(0)):
    from pptx.enum.shapes import MSO_SHAPE

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, W, height)
    _fill(bar, color)
    return bar


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def title_slide(title, subtitle, footer):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, W, H)
    _fill(bg, AZURE)
    _, tf = _tb(s, Emu(700000), Emu(2200000), Emu(10800000), Emu(2000000))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = subtitle
    r2.font.size = Pt(22); r2.font.color.rgb = RGBColor(0xE8, 0xF0, 0xFB)
    _, ff = _tb(s, Emu(700000), Emu(5900000), Emu(10800000), Emu(500000))
    fr = ff.paragraphs[0].add_run(); fr.text = footer
    fr.font.size = Pt(14); fr.font.color.rgb = RGBColor(0xD6, 0xE6, 0xF7)
    return s


def section_slide(kicker, title):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, W, H)
    _fill(bg, DARK)
    _, tf = _tb(s, Emu(700000), Emu(2700000), Emu(10800000), Emu(1600000))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = kicker.upper()
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = AZURE
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = title
    r2.font.size = Pt(40); r2.font.bold = True; r2.font.color.rgb = WHITE
    return s


def content_slide(title, bullets, subtitle=None):
    """bullets: list of (text, level)."""
    s = prs.slides.add_slide(BLANK)
    _band(s)
    _, tt = _tb(s, Emu(600000), Emu(300000), Emu(11000000), Emu(900000))
    tr = tt.paragraphs[0].add_run(); tr.text = title
    tr.font.size = Pt(30); tr.font.bold = True; tr.font.color.rgb = DARK
    if subtitle:
        sp = tt.add_paragraph(); sr = sp.add_run(); sr.text = subtitle
        sr.font.size = Pt(15); sr.font.italic = True; sr.font.color.rgb = GREY

    _, bf = _tb(s, Emu(700000), Emu(1500000), Emu(10800000), Emu(4900000))
    first = True
    for text, level in bullets:
        p = bf.paragraphs[0] if first else bf.add_paragraph()
        first = False
        p.level = level
        run = p.add_run()
        bullet = "•  " if level == 0 else "–  "
        run.text = bullet + text
        run.font.size = Pt(20 - level * 2)
        run.font.color.rgb = DARK if level == 0 else GREY
        p.space_after = Pt(6)
    return s


def code_slide(title, code, caption=None):
    s = prs.slides.add_slide(BLANK)
    _band(s)
    _, tt = _tb(s, Emu(600000), Emu(300000), Emu(11000000), Emu(800000))
    tr = tt.paragraphs[0].add_run(); tr.text = title
    tr.font.size = Pt(28); tr.font.bold = True; tr.font.color.rgb = DARK

    panel = s.shapes.add_shape(1, Emu(600000), Emu(1350000), Emu(11000000), Emu(4650000))
    _fill(panel, CODEBG)
    tf = panel.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Emu(200000); tf.margin_top = Emu(150000)
    first = True
    for line in code.strip("\n").split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = line if line else " "
        r.font.name = "Menlo"; r.font.size = Pt(14); r.font.color.rgb = CODEFG
    if caption:
        _, cf = _tb(s, Emu(600000), Emu(6150000), Emu(11000000), Emu(500000))
        cr = cf.paragraphs[0].add_run(); cr.text = caption
        cr.font.size = Pt(13); cr.font.italic = True; cr.font.color.rgb = GREY
    return s


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

s = title_slide(
    "Agent development got faster",
    "Microsoft Agent Framework · stable orchestration building blocks\n+ Foundry Toolkit for VS Code (GA)",
    "Microsoft Build 2026 Recap — Metro Toronto Azure Community · 25-min session",
)
notes(s, """
Welcome. Two shipping-now things that make agent dev faster:
1) Microsoft Agent Framework now has STABLE orchestration building blocks.
2) Foundry Toolkit for VS Code is GENERALLY AVAILABLE.
We'll cover what they are, then a live demo that uses both. ~25 minutes, code is on GitHub.
""")

s = content_slide("What we'll cover — 25 minutes", [
    ("The problem: agentic apps are mostly orchestration", 0),
    ("Microsoft Agent Framework: successor to Semantic Kernel + AutoGen", 0),
    ("The 5 stable orchestration building blocks (with code)", 0),
    ("Human-in-the-loop: approval gates", 0),
    ("Foundry Toolkit for VS Code (GA): the local dev loop", 0),
    ("Live demo: 'Contoso Outdoors Launch Desk'", 0),
    ("Takeaways + resources", 0),
])
notes(s, "Roadmap slide. ~3 min of framing, ~10 min concepts, ~10 min demo, ~2 min wrap.")

s = content_slide("The real problem isn't the model call", [
    ("A single model call is easy. Real agents are hard because of orchestration:", 0),
    ("Multiple steps, decisions, tools, and several agents coordinating", 1),
    ("Who runs next? In parallel? Who decides we're done?", 1),
    ("When does a human approve a risky action?", 1),
    ("Everyone rebuilds the same coordination plumbing — differently", 0),
    ("Building blocks let you focus on the agents, not the plumbing", 0),
])
notes(s, """
Framing. The hard part of 'agentic' is orchestration complexity, not model access.
If a single tool-calling loop solves it, you don't need a framework. MAF earns its keep
when coordination is the challenge.
""")

s = content_slide("Microsoft Agent Framework (MAF)", [
    ("Open-source SDK for agentic systems — Python and .NET", 0),
    ("The direct successor to Semantic Kernel + AutoGen (same teams)", 0),
    ("AutoGen's simple agent abstractions", 1),
    ("Semantic Kernel's enterprise features: state, type safety, middleware, telemetry", 1),
    ("NEW: graph-based Workflows for explicit multi-agent orchestration", 1),
    ("An Agent = a model + instructions + tools", 0),
    ("pip install agent-framework", 1),
])
notes(s, """
MAF combines AutoGen + SK and adds workflows. Emphasize: not a rewrite of everything —
it's the convergence point. Mention it's a fast-moving 1.x (we pin the version in the repo).
""")

s = code_slide("An agent is small", """
from agent_framework import Agent
from agent_framework.openai import OpenAIChatClient

client = OpenAIChatClient(
    model="gpt-4o-mini",
    azure_endpoint="https://<resource>.openai.azure.com",
    api_version="2024-10-21",
    credential=AzureCliCredential(),   # or api_key=...
)

writer = Agent(client=client, name="Writer",
               instructions="You are a Contoso Outdoors copywriter.")

response = await writer.run("Blurb the Summit backpack.")
print(response.text)
""", caption="demos/00_single_agent.py — the baseline before orchestration")
notes(s, "This is exactly what you first build/test in the Foundry Toolkit Playground & Agent Builder.")

s = content_slide("The 5 orchestration building blocks", [
    ("Sequential — assembly line; each agent builds on the last", 0),
    ("Concurrent — fan-out/fan-in; run in parallel to cut latency", 0),
    ("Handoff — route control to the right specialist", 0),
    ("Group Chat — shared conversation: debate / review / brainstorm", 0),
    ("Magentic — a manager agent dynamically plans & coordinates", 0),
    ("All five support human-in-the-loop (approval + request-info)", 0),
], subtitle="from agent_framework.orchestrations import SequentialBuilder, ConcurrentBuilder, HandoffBuilder, GroupChatBuilder, MagenticBuilder")
notes(s, """
The heart of the talk. These prebuilt patterns handle the coordination boilerplate.
We'll show each with one diagram + a couple lines of code, then demo them.
""")

s = code_slide("1 · Sequential — the pipeline", """
from agent_framework.orchestrations import SequentialBuilder

workflow = SequentialBuilder(
    participants=[outliner, writer, editor],
).build()

result = await workflow.run("Choosing your first 3-season tent")
print(result.get_outputs()[-1])
#   Outliner ─▶ Writer ─▶ Editor
""", caption="demos/01_sequential.py — each agent builds on the previous output")
notes(s, "Content pipeline. Point out how little code the pattern is.")

s = code_slide("2 · Concurrent — the review board", """
from agent_framework.orchestrations import ConcurrentBuilder

workflow = ConcurrentBuilder(
    participants=[seo, legal, brand],   # run in parallel
).build()

result = await workflow.run(marketing_copy)
#        ┌─▶ SEO  ─┐
#        ├─▶ Legal ┤─▶ aggregated feedback
#        └─▶ Brand ┘
""", caption="demos/02_concurrent.py — independent tasks in parallel = lower latency")
notes(s, "Three reviewers look at the same copy at once instead of waiting in line.")

s = code_slide("3 · Handoff — route to a specialist", """
workflow = (
    HandoffBuilder(name="support",
                   participants=[triage, orders, returns, refunds],
                   termination_condition=case_is_closed)
    .with_start_agent(triage)
    .add_handoff(triage, [orders, returns, refunds])
    .add_handoff(returns, [refunds])
    .build()
)
#   triage ─▶ orders / returns ─▶ refunds
""", caption="demos/03_handoff.py — triage transfers control based on context")
notes(s, "Support desk. Triage doesn't answer — it routes. Specialists own their tools.")

s = code_slide("4 · Group Chat + 5 · Magentic", """
# Group chat: shared conversation until convergence
gc = GroupChatBuilder(participants=[copywriter, editor],
                      selection_func=round_robin,
                      termination_condition=shipped).build()

# Magentic: a manager plans & coordinates specialists
mg = MagenticBuilder(participants=[researcher, analyst],
                     manager_agent=launch_manager,
                     max_round_count=8, max_stall_count=2).build()
""", caption="demos/04_group_chat.py · demos/05_magentic.py")
notes(s, """
Group chat = brainstorm/debate; editor says 'SHIP IT' to stop.
Magentic = most autonomous: manager writes a plan, delegates, synthesizes. Bounded by max_* limits.
""")

s = code_slide("Human-in-the-loop = one flag", """
from agent_framework import tool

@tool(approval_mode="always_require")
def publish_blog(title: str, body: str) -> str:
    return f"PUBLISHED: {title}"

# The run pauses and surfaces an approval request:
result = await workflow.run(task)
for ev in result.get_request_info_events():
    req = ev.data                       # function_approval_request
    responses[ev.request_id] = req.to_function_approval_response(approved=True)
result = await workflow.run(responses=responses)
""", caption="demos/06_hitl_approval.py — write tools auto-run, high-impact tools ask first")
notes(s, """
Enterprise story. Read-only tools run automatically; risky/write tools require approval.
Same primitive across all five patterns. This is the trust boundary for autonomous agents.
""")

s = section_slide("Now: how do you build all this?", "Foundry Toolkit for VS Code — now GA")

s = content_slide("Foundry Toolkit for VS Code (GA)", [
    ("The full local AI dev loop — without leaving the editor", 0),
    ("Model Catalog — OpenAI, Anthropic, Google, GitHub, Foundry, ONNX, Ollama", 0),
    ("Playground — chat + params + multimodal, test before you code", 0),
    ("Agent Builder — prompt → production code + MCP tools", 0),
    ("Agent Inspector — debug & visualize agents, set breakpoints", 0),
    ("Evaluation — F1, relevance, coherence, similarity, custom", 0),
    ("Tracing — local OTLP collector; see every step live", 0),
], subtitle="Formerly 'AI Toolkit for VS Code' · install: aka.ms/foundrytk")
notes(s, """
GA matters: supported, stable. It's the on-ramp. You don't provision cloud to start —
discover a model, test in Playground, build in Agent Builder, debug in Agent Inspector,
trace locally, evaluate, then deploy to Foundry Agent Service.
""")

s = content_slide("The inner loop these two enable", [
    ("1. Discover a model in the Model Catalog", 0),
    ("2. Test the prompt in the Playground", 0),
    ("3. Build the agent in Agent Builder → generates code", 0),
    ("4. Compose agents with MAF orchestration building blocks", 0),
    ("5. Debug & visualize in Agent Inspector", 0),
    ("6. Trace locally (OTLP) + Evaluate quality", 0),
    ("7. Deploy hosted agents / workflows to Foundry", 0),
], subtitle="Faster because the loop lives in one place — editor to production")
notes(s, "This is the 'got faster' thesis in one slide. The demo walks this loop.")

s = section_slide("Live demo", "Contoso Outdoors 'Launch Desk'")

s = content_slide("What you'll see in the demo", [
    ("One company, five orchestration patterns, all in VS Code:", 0),
    ("Sequential — blog pipeline (Outliner ▶ Writer ▶ Editor)", 1),
    ("Concurrent — parallel review (SEO ∥ Legal ∥ Brand)", 1),
    ("Handoff — support triage ▶ order / returns / refund", 1),
    ("Magentic — a manager writes a launch brief", 1),
    ("Human-in-the-loop — approve 'publish' before it goes live", 1),
    ("Traces stream into the Foundry Toolkit as it runs", 0),
], subtitle="Repo: agent-framework-foundry-demo · runs on Azure OpenAI (swap backends via .env)")
notes(s, """
DEMO. Follow docs/DEMO_SCRIPT.md. Order: 01 sequential, 02 concurrent, 03 handoff,
05 magentic, then 06 HITL for the 'wow'. Have Agent Inspector / Tracing open in VS Code.
Fallback: set MODEL_BACKEND=github if Azure/Wi-Fi misbehaves. Run `make smoke` beforehand.
""")

s = content_slide("Takeaways", [
    ("Agents are easy; orchestration was the hard part — now it's building blocks", 0),
    ("5 stable patterns: Sequential, Concurrent, Handoff, Group Chat, Magentic", 0),
    ("Human-in-the-loop is a single flag, across every pattern", 0),
    ("Foundry Toolkit (GA) makes the build→debug→trace→evaluate loop local", 0),
    ("Same code runs on Azure OpenAI, Foundry, GitHub Models, or OpenAI", 0),
    ("Start today: pip install agent-framework + install Foundry Toolkit", 0),
])
notes(s, "The one-slide summary if you're short on time. Land the 'got faster' message.")

s = content_slide("Resources", [
    ("Microsoft Agent Framework: learn.microsoft.com/agent-framework", 0),
    ("Orchestrations: learn.microsoft.com/agent-framework/workflows/orchestrations", 0),
    ("Foundry Toolkit for VS Code: code.visualstudio.com/docs/intelligentapps/overview", 0),
    ("Install the toolkit: aka.ms/foundrytk", 0),
    ("This demo repo + slides: <your GitHub URL here>", 0),
    ("Foundry Agent Service: learn.microsoft.com/azure/ai-foundry", 0),
])
notes(s, "Leave this up during Q&A. Replace the repo URL with your GitHub link.")

s = title_slide("Thank you!", "Questions?\nGrab the repo, run `make smoke`, and build an agent today.",
                "Metro Toronto Azure Community · Build 2026 Recap")
notes(s, "Thank the group + organizers. Invite them to clone the repo and try it.")

# ── save ──────────────────────────────────────────────────────────────────────
out = pathlib.Path(__file__).resolve().parent / "Build2026-Recap-Agents.pptx"
prs.save(out)
print(f"Wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
